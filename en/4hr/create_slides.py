#!/usr/bin/env python3
"""
Generate PowerPoint: Agentic RAG Workshop — Pre-Session
Theme: Mahidol University (Navy Blue + Gold)

Output: pre_session_slides.pptx
Upload to Google Slides → File → Import slides → select file

Usage: python3 create_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Mahidol University Theme Colors ─────────────────────────────────────────
NAVY       = RGBColor(0x1A, 0x23, 0x70)   # Deep Mahidol Blue
DARK_NAVY  = RGBColor(0x0F, 0x14, 0x52)   # Darker background
GOLD       = RGBColor(0xC4, 0x99, 0x1C)   # Mahidol Gold
LIGHT_GOLD = RGBColor(0xEA, 0xC7, 0x50)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xDD)
CODE_BG    = RGBColor(0x1E, 0x1E, 0x2E)   # Dark code bg
CODE_GREEN = RGBColor(0x4E, 0xC9, 0x63)   # Green code text
CODE_BLUE  = RGBColor(0x7A, 0xA2, 0xF7)   # Blue keywords
CODE_ORANGE= RGBColor(0xFA, 0xB3, 0x87)   # Orange strings
FAINT_GOLD = RGBColor(0x8B, 0x7A, 0x40)   # Subtle gold for accents


def set_slide_bg(slide, color=NAVY):
    """Set slide background to solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    """Add a textbox and return the text frame."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_run(paragraph, text, size=16, bold=False, color=WHITE,
            font_name="Sarabun", italic=False):
    """Add a styled run to a paragraph."""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return run


def new_para(tf, text="", size=16, bold=False, color=WHITE,
             font_name="Sarabun", align=PP_ALIGN.LEFT, space_after=Pt(6),
             space_before=Pt(0)):
    """Add a new paragraph with a run."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_after = space_after
    p.space_before = space_before
    if text:
        add_run(p, text, size, bold, color, font_name)
    return p


def add_gold_line(slide, top_inches, left=1.0, width=8.0):
    """Add a horizontal gold accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top_inches),
        Inches(width), Pt(3),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_code_block(slide, code_text, left, top, width, height):
    """Add a styled code block with dark background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    # Clear default paragraph
    tf.paragraphs[0].text = ""

    lines = code_text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(1)
        p.space_before = Pt(1)

        # Simple syntax highlighting
        run = p.add_run()
        run.text = line
        run.font.size = Pt(10)
        run.font.name = "Roboto Mono"

        # Color based on content
        stripped = line.strip()
        if stripped.startswith("#"):
            run.font.color.rgb = RGBColor(0x6A, 0x99, 0x55)  # Green comments
        elif any(stripped.startswith(kw) for kw in
                 ["import ", "from ", "for ", "try:", "except ", "if ", "else:", "break", "raise", "def "]):
            run.font.color.rgb = CODE_BLUE
        elif "'" in stripped or '"' in stripped:
            run.font.color.rgb = CODE_ORANGE
        else:
            run.font.color.rgb = CODE_GREEN

    return shape


def add_table(slide, rows_data, left, top, width, height):
    """Add a styled table."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    table = table_shape.table

    # Set column widths proportionally
    total_w = Inches(width)
    col_widths = [total_w // n_cols] * n_cols
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""

            # Set cell background
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY if r % 2 == 0 else RGBColor(0x15, 0x1B, 0x60)

            # Set text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell_text
            run.font.size = Pt(12)
            run.font.name = "Sarabun"
            run.font.bold = (r == 0)
            run.font.color.rgb = GOLD if r == 0 else WHITE

    return table_shape


def add_callout(slide, text, top=4.6):
    """Add a gold callout bar at the bottom."""
    # Background bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(top),
        Inches(9.0), Inches(0.45),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x22, 0x2B, 0x80)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, text, size=13, bold=True, color=LIGHT_GOLD)

    return shape


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    """Cover / Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_NAVY)

    add_gold_line(slide, 1.7)
    add_gold_line(slide, 3.8)

    # Title
    tf = add_textbox(slide, 1.0, 0.8, 8.0, 1.5)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    add_run(p, "Agentic RAG", size=48, bold=True, color=WHITE)
    p2 = new_para(tf, "From Zero to Hero", size=36, bold=True, color=GOLD,
                  space_before=Pt(8))

    # Subtitle
    tf2 = add_textbox(slide, 1.0, 4.0, 8.0, 1.0)
    new_para(tf2, "Workshop 4 Translated English text", size=18, color=LIGHT_GRAY,
             space_after=Pt(4))
    # Remove the default first empty paragraph
    tf2.paragraphs[0].text = ""
    p = tf2.paragraphs[0]
    add_run(p, "Workshop 4 Translated English text", size=18, color=LIGHT_GRAY)
    new_para(tf2, "Translated English text ICT Translated English text", size=18, color=LIGHT_GOLD)


def slide_02_learning(prs):
    """learning outcomes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_gold_line(slide, 0.85, left=0.7, width=3.5)

    tf = add_textbox(slide, 0.7, 0.3, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "learning outcomes", size=32, bold=True, color=GOLD)

    items = [
        ("✂️", "Translated English text RAG (Chunk → Embed → VectorDB)"),
        ("🤖", "Translated English text AI Agent Translated English text Google ADK"),
        ("⭐", "Translated English text Agentic RAG — Agent Translated English text + Translated English text + Translated English text"),
        ("📊", "Translated English text LLM-as-Judge"),
    ]
    tf2 = add_textbox(slide, 1.0, 1.3, 7.5, 3.5)
    tf2.paragraphs[0].text = ""
    for emoji, text in items:
        p = tf2.paragraphs[0] if tf2.paragraphs[0].text == "" else tf2.add_paragraph()
        p.space_after = Pt(18)
        add_run(p, f"{emoji}  ", size=22, color=GOLD)
        add_run(p, text, size=20, color=WHITE)


def slide_03_timeline(prs):
    """Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.85, left=0.7, width=3.0)

    tf = add_textbox(slide, 0.7, 0.3, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "Timeline Translated English text Workshop", size=32, bold=True, color=GOLD)

    table_data = [
        ["Part", "Translated English text", "content"],
        ["📢 Part 1", "1 Translated English text. 20 Translated English text", "Data Pipeline: Chunk → Embed → Qdrant"],
        ["☕ Translated English text", "10 Translated English text", ""],
        ["📢 Part 2", "1 Translated English text. 30 Translated English text", "Agent → Tool → RAG Agent → Agentic RAG"],
        ["🧪 Part 3", "1 Translated English text.", "Translated English text (10 points) + Q&A"],
    ]
    add_table(slide, table_data, 0.7, 1.3, 8.5, 3.0)


def slide_04_register(prs):
    """Translated English text"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.8, left=0.7, width=5.5)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "🌐 Translated English text — myhero.megawiz.co.th", size=28, bold=True, color=GOLD)

    # Left: Steps
    tf2 = add_textbox(slide, 0.7, 1.0, 4.5, 2.8)
    p = tf2.paragraphs[0]
    add_run(p, "Translated English text", size=16, bold=True, color=LIGHT_GOLD)

    steps = [
        "1. Translated English text myhero.megawiz.co.th/student-portal",
        "2. Translated English text OTP",
        "3. Translated English text OTP → Translated English text",
        "4. Translated English text-Translated English textstudents",
    ]
    for s in steps:
        new_para(tf2, s, size=14, color=WHITE, space_after=Pt(6))

    # Right: Benefits
    tf3 = add_textbox(slide, 5.3, 1.0, 4.2, 2.8)
    p = tf3.paragraphs[0]
    add_run(p, "Translated English text", size=16, bold=True, color=LIGHT_GOLD)

    benefits = [
        "📄 Translated English text (Notebook)",
        "📝 Translated English text + Translated English text",
        "🏆 Certificate Translated English text workshop",
    ]
    for b in benefits:
        new_para(tf3, b, size=14, color=WHITE, space_after=Pt(8))

    add_callout(slide, "🎯 Translated English text workshop Translated English text")


def slide_05_prep(prs):
    """Translated English text"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.85, left=0.7, width=5.0)

    tf = add_textbox(slide, 0.7, 0.3, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "⚠️ Translated English text", size=32, bold=True, color=GOLD)

    items = [
        ("✅", "Translated English text myhero.megawiz.co.th/student-portal"),
        ("🔑", "Translated English text Gemini API Key"),
        ("💻", "Translated English text Google Colab Translated English text"),
        ("🧪", "Translated English text API Key Translated English text"),
    ]
    tf2 = add_textbox(slide, 1.0, 1.4, 8.0, 3.0)
    tf2.paragraphs[0].text = ""
    for i, (emoji, text) in enumerate(items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(18)
        add_run(p, f"{i+1}. {emoji}  ", size=22, bold=True, color=LIGHT_GOLD)
        add_run(p, text, size=20, color=WHITE)


def slide_06_apikey(prs):
    """Translated English text Gemini API Key"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.8, left=0.7, width=4.5)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "🔑 Translated English text Gemini API Key", size=28, bold=True, color=GOLD)

    steps = [
        'Translated English text aistudio.google.com',
        'Login Translated English text Google Account',
        'Translated English text "Get API Key" (Translated English text)',
        'Translated English text "Create API Key"',
        'Translated English text Project → Create API key in new project',
        'Copy API Key Translated English text (Translated English text!)',
    ]
    tf2 = add_textbox(slide, 0.7, 1.0, 8.5, 3.2)
    tf2.paragraphs[0].text = ""
    for i, step in enumerate(steps):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(10)
        add_run(p, f"{i+1}. ", size=16, bold=True, color=LIGHT_GOLD)
        add_run(p, step, size=16, color=WHITE)

    add_callout(slide, "⚠️ Free tier: 15 RPM (requests/minute) — Translated English text workshop")


def slide_07_secrets(prs):
    """Translated English text API Key Translated English text Colab Secrets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.8, left=0.7, width=5.5)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "🔒 Translated English text API Key Translated English text Colab Secrets", size=28, bold=True, color=GOLD)

    steps = [
        'Translated English text Google Colab',
        'Translated English text 🔑 Translated English text (Translated English text) → "Secrets"',
        'Translated English text "Add new secret"',
        'Name: GOOGLE_API_KEY',
        'Value: Translated English text API Key Translated English text copy Translated English text',
        'Translated English text toggle "Notebook access" ✅',
    ]
    tf2 = add_textbox(slide, 0.7, 1.0, 8.5, 3.2)
    tf2.paragraphs[0].text = ""
    for i, step in enumerate(steps):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(10)
        add_run(p, f"{i+1}. ", size=16, bold=True, color=LIGHT_GOLD)
        add_run(p, step, size=16, color=WHITE)

    add_callout(slide, "🎯 Translated English text — Translated English text notebook")


def slide_08_ratelimit(prs):
    """Error 429"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.8, left=0.7, width=4.0)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "🚨 Error 429 — Rate Limit Translated English text?", size=28, bold=True, color=GOLD)

    code = "google.api_core.exceptions.ResourceExhausted: 429\nYou exceeded your current quota"
    add_code_block(slide, code, 0.7, 1.0, 8.5, 0.7)

    tf2 = add_textbox(slide, 0.7, 1.9, 8.5, 0.5)
    p = tf2.paragraphs[0]
    add_run(p, "Translated English text: ", size=16, bold=True, color=GOLD)
    add_run(p, "Translated English text request Translated English text/Translated English text", size=16, color=WHITE)

    # Free tier limits table
    tf3 = add_textbox(slide, 0.7, 2.5, 8.5, 0.4)
    p = tf3.paragraphs[0]
    add_run(p, "Free Tier Limits (Gemini 2.5 Flash):", size=14, bold=True, color=LIGHT_GOLD)

    table_data = [
        ["Limit", "Translated English text"],
        ["RPM (Requests/min)", "15"],
        ["TPM (Tokens/min)", "1,000,000"],
        ["RPD (Requests/day)", "1,500"],
    ]
    add_table(slide, table_data, 0.7, 3.0, 8.5, 2.0)


def slide_09_fix429(prs):
    """Translated English text Error 429"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.8, left=0.7, width=3.0)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "💡 Translated English text Error 429", size=28, bold=True, color=GOLD)

    # Quick fix section
    tf2 = add_textbox(slide, 0.7, 0.9, 4.0, 1.5)
    p = tf2.paragraphs[0]
    add_run(p, "Translated English text:", size=16, bold=True, color=LIGHT_GOLD)
    p2 = new_para(tf2, "1. Translated English text 1-2 Translated English text run Translated English text", size=13, color=WHITE, space_after=Pt(6))
    p3 = new_para(tf2, "2. Translated English text run cell Translated English text — Translated English text output Translated English text", size=13, color=WHITE)

    # Permanent fix label
    tf3 = add_textbox(slide, 5.0, 0.9, 4.5, 0.4)
    p = tf3.paragraphs[0]
    add_run(p, "Translated English text (Translated English text):", size=16, bold=True, color=LIGHT_GOLD)

    # Code block
    code = (
        "import time\n"
        "for attempt in range(3):\n"
        "    try:\n"
        "        response = client.models.generate_content(...)\n"
        "        break\n"
        "    except Exception as e:\n"
        "        if '429' in str(e):\n"
        "            print(f'⏳ Rate limit, Translated English text {30*(attempt+1)} Translated English text...')\n"
        "            time.sleep(30 * (attempt + 1))\n"
        "        else:\n"
        "            raise"
    )
    add_code_block(slide, code, 0.5, 2.5, 9.0, 2.5)


def slide_10_test(prs):
    """Translated English text"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.8, left=0.7, width=3.5)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "🧪 Translated English text", size=28, bold=True, color=GOLD)

    tf2 = add_textbox(slide, 0.7, 0.75, 8.5, 0.3)
    p = tf2.paragraphs[0]
    add_run(p, "Translated English text Colab → Translated English text notebook Translated English text → Run:", size=13, color=LIGHT_GRAY)

    code = (
        "# 1. Translated English text API Key\n"
        "from google.colab import userdata\n"
        "import os\n"
        "os.environ['GOOGLE_API_KEY'] = userdata.get('GOOGLE_API_KEY')\n"
        "\n"
        "# 2. Translated English text\n"
        "from google import genai\n"
        "client = genai.Client(api_key=os.environ['GOOGLE_API_KEY'])\n"
        "resp = client.models.generate_content(\n"
        "    model='gemini-2.5-flash',\n"
        "    contents='Translated English text 1 Translated English text'\n"
        ")\n"
        "print(f'✅ Translated English text: {resp.text}')"
    )
    add_code_block(slide, code, 0.7, 1.1, 8.5, 2.8)

    add_callout(slide, "Translated English text ✅ = Translated English text!", top=4.2)


def slide_11_checklist(prs):
    """Checklist"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.85, left=0.7, width=3.5)

    tf = add_textbox(slide, 0.7, 0.3, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "📋 Checklist Translated English text", size=32, bold=True, color=GOLD)

    items = [
        "Translated English text myhero.megawiz.co.th/student-portal Translated English text",
        "Translated English text Google Account",
        "Translated English text Gemini API Key Translated English text",
        "Translated English text Key Translated English text Colab Secrets Translated English text",
        "Run Translated English text → Translated English text ✅",
    ]
    tf2 = add_textbox(slide, 1.0, 1.3, 8.0, 3.0)
    tf2.paragraphs[0].text = ""
    for i, item in enumerate(items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(14)
        add_run(p, "☐  ", size=22, bold=True, color=LIGHT_GOLD)
        add_run(p, item, size=18, color=WHITE)

    add_callout(slide, "🎯 Translated English text checklist Translated English text → Translated English text workshop Translated English text!")


def slide_12_cert(prs):
    """Certificate"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.85, left=0.7, width=2.5)

    tf = add_textbox(slide, 0.7, 0.3, 8.5, 0.6)
    p = tf.paragraphs[0]
    add_run(p, "🏆 Certificate", size=32, bold=True, color=GOLD)

    tf2 = add_textbox(slide, 0.7, 1.2, 8.5, 0.4)
    p = tf2.paragraphs[0]
    add_run(p, "Translated English text workshop Translated English text:", size=16, color=LIGHT_GRAY)

    items = [
        ("🤖", "Translated English textpointsTranslated English text (AI Grading)"),
        ("📜", "Translated English text Certificate Translated English text myhero.megawiz.co.th"),
        ("📥", "Translated English text (PDF)"),
    ]
    tf3 = add_textbox(slide, 1.0, 1.8, 8.0, 2.5)
    tf3.paragraphs[0].text = ""
    for i, (emoji, text) in enumerate(items):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.space_after = Pt(18)
        add_run(p, f"{emoji}  ", size=24, color=GOLD)
        add_run(p, text, size=20, color=WHITE)


def slide_13_help(prs):
    """Translated English text?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_line(slide, 0.75, left=0.7, width=2.0)

    tf = add_textbox(slide, 0.7, 0.2, 8.5, 0.5)
    p = tf.paragraphs[0]
    add_run(p, "❓ Translated English text?", size=32, bold=True, color=GOLD)

    table_data = [
        ["Translated English text", "Translated English text"],
        ["Translated English text myhero Translated English text", "Translated English text spam / Translated English text"],
        ["Translated English text OTP", "Translated English text spam folder / Translated English text 1 Translated English text"],
        ["API key not valid", "Translated English text key Translated English text aistudio.google.com"],
        ["Error 429", "Translated English text 1-2 Translated English text"],
        ["Colab Translated English text", "Runtime → Change runtime type → T4 GPU"],
    ]
    add_table(slide, table_data, 0.7, 0.9, 8.5, 3.2)

    add_callout(slide, "Translated English text / TA Translated English text 🙋", top=4.4)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # Build all 13 slides
    slide_01_cover(prs)
    slide_02_learning(prs)
    slide_03_timeline(prs)
    slide_04_register(prs)
    slide_05_prep(prs)
    slide_06_apikey(prs)
    slide_07_secrets(prs)
    slide_08_ratelimit(prs)
    slide_09_fix429(prs)
    slide_10_test(prs)
    slide_11_checklist(prs)
    slide_12_cert(prs)
    slide_13_help(prs)

    output_path = "pre_session_slides.pptx"
    prs.save(output_path)
    print(f"✅ Translated English text: {output_path}")
    print(f"📤 Upload Translated English text Google Slides:")
    print(f"   1. Translated English text slides.google.com")
    print(f"   2. Translated English text + (Blank) → File → Import slides")
    print(f"   3. Translated English text upload Translated English text .pptx Translated English text Google Drive Translated English text Google Slides")


if __name__ == "__main__":
    main()
